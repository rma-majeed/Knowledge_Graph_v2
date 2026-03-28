"""Generate minimal synthetic fixtures for testing."""
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

FIXTURES = Path(__file__).parent


def make_sample_pdf():
    """Create a 2-page PDF with text and a table."""
    doc = fitz.open()

    # Page 1 — plain text
    page1 = doc.new_page(width=595, height=842)  # A4
    page1.insert_text(
        (72, 72),
        "Automotive Consulting Report\n\nExecutive Summary\n\n"
        "This document covers EV supply chain strategy for OEM clients. "
        "Toyota and BMW are the primary stakeholders. The recommendation is "
        "to diversify battery suppliers beyond CATL to include Panasonic and BYD.\n\n"
        "Section 1: Market Analysis\n\n"
        "The global EV market is projected to reach $800B by 2030. "
        "Key technology trends include solid-state batteries and 800V charging architecture.",
        fontsize=11,
    )

    # Page 2 — text + simulate table content
    page2 = doc.new_page(width=595, height=842)
    page2.insert_text(
        (72, 72),
        "Section 2: Supplier Recommendations\n\n"
        "OEM\tPrimary Supplier\tAlternative\n"
        "Toyota\tPanasonic\tCATL\n"
        "BMW\tSamsung SDI\tBYD\n\n"
        "Conclusion: Dual-source battery strategy reduces supply chain risk by 40%.",
        fontsize=11,
    )

    doc.save(FIXTURES / "sample.pdf")
    doc.close()
    print(f"Created: {FIXTURES / 'sample.pdf'}")


def make_sample_pptx():
    """Create a 3-slide PPTX with text, speaker notes, and a table."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1 — title + body text
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide1.shapes.title.text = "EV Strategy Overview"
    slide1.placeholders[1].text = (
        "Key findings:\n"
        "- Toyota leads in hybrid technology\n"
        "- Battery costs declining 15% per year\n"
        "- 800V architecture is next inflection point"
    )
    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = "Speaker note: Emphasize cost trajectory to CFO audience."

    # Slide 2 — blank with text box
    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    txBox.text_frame.text = (
        "Market share analysis\n\n"
        "BYD has captured 35% of China EV market. "
        "Tesla maintains 18% global share. "
        "Legacy OEMs collectively hold 47% share."
    )
    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = "Speaker note: Updated Q3 2025 data — verify with client before sharing."

    # Slide 3 — table
    slide3 = prs.slides.add_slide(blank_layout)
    rows, cols = 3, 3
    table = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(8), Inches(3)).table
    headers = ["OEM", "EV Sales 2024", "Battery Partner"]
    data = [
        ["Toyota", "450,000", "Panasonic / Prime Planet"],
        ["BMW", "375,000", "Samsung SDI"],
    ]
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header
    for row_idx, row in enumerate(data):
        for col_idx, val in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = val

    prs.save(FIXTURES / "sample.pptx")
    print(f"Created: {FIXTURES / 'sample.pptx'}")


if __name__ == "__main__":
    make_sample_pdf()
    make_sample_pptx()
    print("Fixtures ready.")
