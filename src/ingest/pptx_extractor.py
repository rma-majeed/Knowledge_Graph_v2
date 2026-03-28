"""PPTX text extraction using python-pptx.

Extracts text from slide shapes, table cells, and speaker notes.
Returns a list of slide dicts with 1-indexed slide numbers.

Usage:
    from src.ingest.pptx_extractor import extract_pptx
    slides = extract_pptx("path/to/presentation.pptx")
    # slides[0] == {"slide_num": 1, "text": "..."}
"""
from __future__ import annotations

from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.shapes.base import BaseShape


def _extract_shape_text(shape: BaseShape) -> str:
    """Extract text from a single shape, including table cells.

    Args:
        shape: A python-pptx BaseShape object.

    Returns:
        Text string extracted from the shape, or empty string if shape has no text.
    """
    parts: list[str] = []

    if shape.has_table:
        # Table: extract each cell row by row, tab-separated
        table = shape.table
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    elif hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            parts.append(text)

    return "\n".join(parts)


def extract_pptx(filepath: Union[str, Path]) -> list[dict]:
    """Extract text from all slides of a PPTX file.

    For each slide, extracts:
    - Text from all shapes (titles, body text, text boxes, bullet points)
    - Table cell content (row by row, cells tab-separated)
    - Speaker notes (appended as "[NOTES] {notes_text}")

    Args:
        filepath: Path to the PPTX file (str or pathlib.Path).

    Returns:
        List of dicts, one per slide, in presentation order:
        [
            {"slide_num": 1, "text": "..."},  # 1-indexed
            {"slide_num": 2, "text": "..."},
            ...
        ]

    Raises:
        FileNotFoundError: If the file does not exist.
    """
    filepath = Path(filepath)
    if not filepath.exists():
        raise FileNotFoundError(f"PPTX not found: {filepath}")

    prs = Presentation(str(filepath))
    slides: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        parts: list[str] = []

        # Extract text from all shapes on the slide
        for shape in slide.shapes:
            shape_text = _extract_shape_text(shape)
            if shape_text:
                parts.append(shape_text)

        # Extract speaker notes
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[NOTES] {notes_text}")
            except Exception:
                # Notes frame can be malformed in some PPTXs — degrade gracefully
                pass

        combined_text = "\n".join(p for p in parts if p)

        slides.append({
            "slide_num": slide_idx + 1,  # Convert 0-indexed to 1-indexed
            "text": combined_text,
        })

    return slides
